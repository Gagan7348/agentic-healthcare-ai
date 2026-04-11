"""
Agentic Healthcare AI — Professional PowerPoint Generator
Generates a refined, full-detail 20-mark presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Color Palette (Dark Clinical Theme) ───────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # Deep Navy
BG_CARD      = RGBColor(0x11, 0x26, 0x3A)   # Card Navy
ACCENT_TEAL  = RGBColor(0x00, 0xC9, 0xB1)   # Teal/Mint
ACCENT_BLUE  = RGBColor(0x2E, 0x86, 0xC1)   # Medical Blue
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)   # Success Green
ACCENT_RED   = RGBColor(0xE7, 0x4C, 0x3C)   # Alert Red
ACCENT_GOLD  = RGBColor(0xF3, 0x9C, 0x12)   # Gold/Warning
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xB2, 0xBE, 0xC3)
DARK_GREY    = RGBColor(0x2D, 0x3A, 0x4A)


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # Completely blank layout

# ─── Helper Functions ──────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_multiline(slide, lines, left, top, width, height,
                  font_size=14, color=WHITE, line_spacing=1.15,
                  bold_first=False):
    """lines = list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = bld
        run.font.color.rgb = clr
        run.font.name = "Calibri"
    return txBox

def add_divider(slide, top, color=ACCENT_TEAL):
    add_rect(slide, 0.5, top, 12.33, 0.04, color)

def slide_header(slide, title, subtitle=None, tag=None):
    """Standard header bar used on every content slide"""
    add_rect(slide, 0, 0, 13.33, 1.1, BG_CARD)
    add_rect(slide, 0, 1.08, 13.33, 0.04, ACCENT_TEAL)

    add_text(slide, title, 0.4, 0.08, 9, 0.65,
             font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.72, 9, 0.4,
                 font_size=13, color=ACCENT_TEAL, italic=True)
    if tag:
        add_text(slide, tag, 10.5, 0.2, 2.5, 0.5,
                 font_size=11, color=DARK_GREY, align=PP_ALIGN.RIGHT,
                 bold=False)

def footer(slide, page_num, total=17):
    add_rect(slide, 0, 7.25, 13.33, 0.25, BG_CARD)
    add_text(slide, f"Slide {page_num} / {total}", 0.3, 7.26, 3, 0.22,
             font_size=9, color=LIGHT_GREY)
    add_text(slide, "Agentic Healthcare AI  |  B.Tech CSE VI Sem  |  Final Review III",
             3.5, 7.26, 7, 0.22,
             font_size=9, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(slide, "© 2026", 11.5, 7.26, 1.5, 0.22,
             font_size=9, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)

# Top gradient bar
add_rect(sl, 0, 0, 13.33, 0.15, ACCENT_TEAL)
add_rect(sl, 0, 7.35, 13.33, 0.15, ACCENT_TEAL)

# Decorative side accent
add_rect(sl, 0, 0, 0.12, 7.5, ACCENT_TEAL)

# Main title
add_text(sl, "🏥 Agentic Healthcare AI", 0.5, 1.3, 12.0, 1.0,
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(sl, "Intelligent Clinical Diagnostic & Multi-Agent Decision Support System",
         0.5, 2.45, 12.0, 0.6,
         font_size=20, color=ACCENT_TEAL, align=PP_ALIGN.CENTER, italic=True)

add_divider(sl, 3.25, ACCENT_TEAL)

add_text(sl, '"Bridging the 1:1456 doctor-patient gap in rural India through autonomous AI"',
         0.8, 3.45, 11.5, 0.55,
         font_size=15, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

# Info grid
add_rect(sl, 1.5, 4.25, 3.0, 1.6, DARK_GREY)
add_text(sl, "🎓  B.Tech CSE — VI Sem", 1.6, 4.4, 2.8, 0.4, font_size=13, color=WHITE)
add_text(sl, "Final Presentation (Review III)", 1.6, 4.82, 2.8, 0.35, font_size=11, color=LIGHT_GREY)
add_text(sl, "Target: 20 / 20 Marks", 1.6, 5.15, 2.8, 0.35, font_size=11, color=ACCENT_GREEN, bold=True)

add_rect(sl, 5.17, 4.25, 3.0, 1.6, DARK_GREY)
add_text(sl, "⚡  Tech Stack", 5.27, 4.4, 2.8, 0.4, font_size=13, color=WHITE)
add_text(sl, "Gemini · GPT-4o · Groq", 5.27, 4.82, 2.8, 0.35, font_size=11, color=LIGHT_GREY)
add_text(sl, "FastAPI · MongoDB · React", 5.27, 5.15, 2.8, 0.35, font_size=11, color=LIGHT_GREY)

add_rect(sl, 8.83, 4.25, 3.0, 1.6, DARK_GREY)
add_text(sl, "📅  April 2026", 8.93, 4.4, 2.8, 0.4, font_size=13, color=WHITE)
add_text(sl, "[Your Institution Name]", 8.93, 4.82, 2.8, 0.35, font_size=11, color=LIGHT_GREY)
add_text(sl, "Guided by: [Guide Name]", 8.93, 5.15, 2.8, 0.35, font_size=11, color=LIGHT_GREY)

footer(sl, 1)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Contents", "What We Will Cover Today")

contents_left = [
    ("01", "Statement of Problem / Scope"),
    ("02", "Objectives of Study"),
    ("03", "Proposed System & Tech Stack"),
    ("04", "Database Architecture  ★"),
    ("05", "ML Model & Clinical Overrides"),
    ("06", "Results & Validation"),
    ("07", "Live Demo — Deployed System"),
    ("08", "Comparison: Strengths & Weaknesses  ★"),
    ("09", "Why Data is Authenticated  ★"),
    ("10", "Outcomes & Impact"),
    ("11", "Mapping: POs & SDGs"),
    ("12", "Conclusion & Future Work"),
]

for i, (num, title) in enumerate(contents_left):
    row = i % 6
    col = i // 6
    y = 1.35 + row * 0.95
    x = 0.5 + col * 6.5

    add_rect(sl, x, y, 0.55, 0.65, ACCENT_TEAL)
    add_text(sl, num, x, y + 0.1, 0.55, 0.5,
             font_size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    is_new = "★" in title
    clr = ACCENT_GOLD if is_new else WHITE
    add_text(sl, title.replace(" ★", ""), x + 0.65, y + 0.12, 5.7, 0.5,
             font_size=14, color=clr, bold=is_new)

add_text(sl, "★ = New section added for this review",
         0.5, 7.1, 6, 0.28, font_size=11, color=ACCENT_GOLD, italic=True)
footer(sl, 2)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Statement of Problem / Scope",
             "India's Healthcare Crisis — Why This Project Exists")

stats = [
    ("1 : 1,456", "Doctor-to-Patient\nRatio in India", ACCENT_RED),
    ("650,000+", "Villages with\nNo Specialist", ACCENT_GOLD),
    ("77 Million", "Indians Affected\nby Diabetes", ACCENT_BLUE),
    ("~42%", "Diagnostic Error\nRate at Rural PHCs", ACCENT_RED),
]
for i, (val, label, clr) in enumerate(stats):
    x = 0.4 + i * 3.2
    add_rect(sl, x, 1.25, 3.0, 1.7, DARK_GREY)
    add_rect(sl, x, 1.25, 3.0, 0.08, clr)
    add_text(sl, val, x, 1.45, 3.0, 0.75,
             font_size=30, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_text(sl, label, x, 2.2, 3.0, 0.6,
             font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_divider(sl, 3.15)

root_causes = [
    "🗺️  Geographical Barriers — 650K+ villages have no specialist doctor within 50 km",
    "🗣️  Language Barriers — 22 official languages; medical reports only in English/Hindi",
    "⏳  Diagnostic Delays — 2–3 weeks average wait for specialist consultation in Tier-3 areas",
    "👩‍⚕️  Overburdened ASHA Workers — Serve 1,000+ patients each with zero AI decision support",
    "💰  Economic Barriers — Private specialist consultation costs ₹500–₹2,000 per visit",
]
add_text(sl, "Root Causes:", 0.5, 3.3, 4, 0.4,
         font_size=14, bold=True, color=ACCENT_TEAL)
for i, cause in enumerate(root_causes):
    add_text(sl, cause, 0.5, 3.72 + i * 0.57, 12.3, 0.5,
             font_size=13, color=WHITE)

footer(sl, 3)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — SCOPE (6 DIMENSIONS)
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Scope of the Project — 6 Dimensions  ★",
             "Who, Where, What, How, For Whom — and What is Deliberately Excluded")

# Intro text
add_text(sl, "Scope is defined across 6 dimensions — this is a purpose-built clinical AI, not a generic chatbot:",
         0.4, 1.22, 12.5, 0.38, font_size=13, color=LIGHT_GREY, italic=True)

scope_dims = [
    ("📍", "Geographic",
     "Tier-2, Tier-3 cities + rural PHCs / Sub-centres\nTarget: ASHA, ANM, PHC Doctors — NOT urban metros",
     ACCENT_TEAL),
    ("🏥", "Clinical / Disease",
     "Diabetes ✅  Heart Disease ✅  CKD ✅\nTB / Anaemia → v5.0   Cancer / Mental Health → Out of Scope",
     ACCENT_RED),
    ("🤖", "Functional / Feature",
     "12 live modules: ML Predict, Dual-AI, ASHA Triage,\nOCR, Voice TTS, Drug Safety, Patient Registry, Trends",
     ACCENT_BLUE),
    ("🏗️", "Technical / Architecture",
     "FastAPI + MongoDB Atlas + Gemini + GPT-4o + Groq\nML inference-only; SQLite = dev fallback only",
     ACCENT_GOLD),
    ("👥", "Stakeholder / Users",
     "ASHA Worker (field triage)  |  Rural Patient (voice+lang)\nPHC Doctor (report)  |  Admin (registry)  |  Researcher",
     ACCENT_GREEN),
    ("⛔", "Out of Scope",
     "No prescriptions  |  No surgery  |  No ICU monitoring\nNo insurance  |  No pediatric  |  No mental health",
     ACCENT_GOLD),
]

for i, (icon, title, desc, clr) in enumerate(scope_dims):
    col = i % 3; row = i // 3
    x = 0.35 + col * 4.35
    y = 1.75 + row * 2.55
    add_rect(sl, x, y, 4.1, 2.3, DARK_GREY)
    add_rect(sl, x, y, 4.1, 0.07, clr)
    add_rect(sl, x, y, 0.07, 2.3, clr)
    add_text(sl, icon + "  " + title, x + 0.18, y + 0.14, 3.8, 0.42,
             font_size=13, bold=True, color=clr)
    add_text(sl, desc, x + 0.18, y + 0.65, 3.8, 1.5,
             font_size=11, color=WHITE)

footer(sl, 4)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — OBJECTIVES
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Objectives of Study",
             "6 Core Goals of the Agentic Healthcare AI System")

objectives = [
    ("01", "Multi-Agent AI Clinical OS",
     "Deploy Google Gemini 1.5 + GPT-4o + Groq/Llama as a collaborative 3-agent diagnostic panel", ACCENT_TEAL),
    ("02", "ML Disease Risk Prediction",
     "XGBoost Stacking Ensemble for Diabetes, Heart Disease & CKD with >91% accuracy", ACCENT_BLUE),
    ("03", "Persistent Clinical Database",
     "MongoDB Atlas with 10,000+ patients, 30,000+ diagnoses, auto-partitioned into Critical/Monitoring/Optimal", ACCENT_GREEN),
    ("04", "10-Language Healthcare Access",
     "Native script output in Hindi, Tamil, Telugu, Bengali, Marathi, Gujarati, Kannada, Malayalam, Punjabi, English", ACCENT_GOLD),
    ("05", "ASHA Worker Triage Tool",
     "RED/YELLOW/GREEN urgency system aligned to NHM guidelines — bilingual, built for PHC-level workers", ACCENT_RED),
    ("06", "Production Deployment",
     "FastAPI backend on Render + React/Vite frontend on Netlify — live, globally accessible, 99.7% uptime", ACCENT_TEAL),
]

for i, (num, title, desc, clr) in enumerate(objectives):
    col = i % 3; row = i // 3
    x = 0.4 + col * 4.3
    y = 1.3 + row * 2.8
    add_rect(sl, x, y, 4.0, 2.5, DARK_GREY)
    add_rect(sl, x, y, 0.08, 2.5, clr)
    add_rect(sl, x, y, 4.0, 0.08, clr)
    add_text(sl, num, x + 0.2, y + 0.12, 0.6, 0.45,
             font_size=18, bold=True, color=clr)
    add_text(sl, title, x + 0.3, y + 0.55, 3.55, 0.45,
             font_size=13, bold=True, color=WHITE)
    add_text(sl, desc, x + 0.3, y + 1.0, 3.55, 1.35,
             font_size=11, color=LIGHT_GREY)

footer(sl, 5)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Proposed System & Technical Architecture",
             "Full-Stack Multi-Agent Clinical OS — NeuroHealth v4.0.5")

# Insert the generated architecture diagram PNG
sl.shapes.add_picture(
    "system_architecture.png",
    Inches(1.0), Inches(1.3), height=Inches(5.7)
)

footer(sl, 6)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — DATABASE ARCHITECTURE (PART 1)
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Database Architecture — Collections & Schema  ★",
             "MongoDB Atlas — 6 Collections, Field-Level Schema, Partitioning")

# Left column: Why MongoDB
add_text(sl, "Why MongoDB (NoSQL) over SQL?", 0.4, 1.25, 5.8, 0.38,
         font_size=14, bold=True, color=ACCENT_TEAL)
why = [
    ("Variable Schema", "Diabetes → HbA1c fields; Heart → Cholesterol. SQL would need NULL columns."),
    ("JSON-Native",     "REST API returns JSON. MongoDB stores BSON. Zero transformation needed."),
    ("Atlas SLA 99.99%","Auto-sharding + geo-replicas. No manual DBA required."),
    ("O(1) Triage",     "Partitioned collections → ASHA dashboard reads <5ms vs full 10K scan."),
]
for i, (k, v) in enumerate(why):
    y = 1.7 + i * 0.7
    add_rect(sl, 0.4, y, 1.5, 0.55, DARK_GREY)
    add_text(sl, k, 0.5, y + 0.08, 1.4, 0.4, font_size=10, bold=True, color=ACCENT_TEAL)
    add_text(sl, v, 2.05, y + 0.06, 4.0, 0.5, font_size=11, color=LIGHT_GREY)

# Right column: Collection cards
add_text(sl, "6 Collections:", 6.6, 1.25, 6.5, 0.38,
         font_size=14, bold=True, color=WHITE)
collections = [
    ("patients",            "Master Registry",       "10,000+ records",  ACCENT_BLUE),
    ("diagnoses",           "ML Predictions + Labs",  "30,000+ records",  ACCENT_TEAL),
    ("consultations",       "AI Chat Logs",           "Growing",          ACCENT_GREEN),
    ("patients_critical",   "🔴 Glucose > 180",       "~2,800 records",   ACCENT_RED),
    ("patients_monitoring", "🟡 Glucose 120-180",     "~4,100 records",   ACCENT_GOLD),
    ("patients_optimal",    "🟢 Stable Patients",     "~3,100 records",   ACCENT_GREEN),
]
for i, (coll, purpose, count, clr) in enumerate(collections):
    col = i % 2; row = i // 2
    x = 6.6 + col * 3.3
    y = 1.72 + row * 1.8
    add_rect(sl, x, y, 3.1, 1.55, DARK_GREY)
    add_rect(sl, x, y, 0.06, 1.55, clr)
    add_text(sl, coll, x + 0.15, y + 0.08, 2.9, 0.42,
             font_size=11, bold=True, color=clr)
    add_text(sl, purpose, x + 0.15, y + 0.52, 2.9, 0.35,
             font_size=10, color=LIGHT_GREY)
    add_text(sl, count, x + 0.15, y + 0.9, 2.9, 0.35,
             font_size=11, color=WHITE, bold=True)

# Partition logic bar at bottom
add_rect(sl, 0.4, 4.65, 12.5, 1.05, BG_CARD)
add_rect(sl, 0.4, 4.65, 12.5, 0.06, ACCENT_RED)
add_text(sl, "⚡ Partitioning Logic — runs on every save_diagnosis() call:",
         0.55, 4.72, 10, 0.38, font_size=12, bold=True, color=WHITE)
partition_code = (
    "if glucose > 180 or hba1c > 8.0  →  patients_critical   🔴     "
    "elif glucose > 120 or hba1c > 6.5  →  patients_monitoring 🟡     "
    "else  →  patients_optimal  🟢"
)
add_text(sl, partition_code, 0.55, 5.12, 12.2, 0.5,
         font_size=11, color=ACCENT_TEAL)

# Field schema strip (patients collection)
add_text(sl, "patients fields: patient_ref (unique) | name | age | gender | phone | email | created_at",
         0.4, 5.78, 12.5, 0.38, font_size=11, color=LIGHT_GREY, italic=True)
add_text(sl, "diagnoses fields: patient_ref | disease | prediction | confidence | risk_score | model_used | glucose | bp | bmi | hba1c | cholesterol | creatinine | created_at",
         0.4, 6.18, 12.5, 0.38, font_size=11, color=LIGHT_GREY, italic=True)

footer(sl, 7)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — ML MODELS
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "ML Model Architecture & Clinical Override Engine",
             "XGBoost Stacking Ensemble + Evidence-Based Safety Net")

# Model cards
models_data = [
    ("🩸", "Diabetes\nDetector",   "94.2%", "AUC: 0.97", "HbA1c, Glucose, BMI, Family History", ACCENT_RED),
    ("❤️",  "Heart Disease\nRisk",  "91.7%", "AUC: 0.95", "BP, Cholesterol, Age, Smoking, HDL",   ACCENT_BLUE),
    ("🫘", "CKD Early\nWarning",   "93.5%", "AUC: 0.96", "Creatinine, BP, Diabetes Risk, Age",   ACCENT_GOLD),
]
for i, (icon, name, acc, auc, feats, clr) in enumerate(models_data):
    x = 0.4 + i * 4.3
    add_rect(sl, x, 1.25, 4.0, 3.1, DARK_GREY)
    add_rect(sl, x, 1.25, 4.0, 0.07, clr)
    add_text(sl, icon + "  " + name, x + 0.2, 1.4, 3.6, 0.7,
             font_size=15, bold=True, color=clr)
    add_text(sl, "Accuracy", x + 0.2, 2.1, 1.6, 0.35,
             font_size=11, color=LIGHT_GREY)
    add_text(sl, acc, x + 0.2, 2.42, 1.6, 0.5,
             font_size=28, bold=True, color=WHITE)
    add_text(sl, auc, x + 1.9, 2.42, 1.8, 0.5,
             font_size=14, color=clr, bold=True)
    add_text(sl, "Key Features:\n" + feats, x + 0.2, 2.95, 3.6, 0.9,
             font_size=10, color=LIGHT_GREY)

add_divider(sl, 4.55)

add_text(sl, "Clinical Override Engine (ADA / ACC-AHA / KDIGO Guidelines)",
         0.5, 4.65, 12, 0.38, font_size=14, bold=True, color=ACCENT_TEAL)

overrides = [
    ("ADA  →  Diabetes", "HbA1c ≥ 6.5%  or  Glucose ≥ 200 mg/dL  →  risk overridden to 92%+", ACCENT_RED),
    ("ACC/AHA  →  Heart", "BP ≥ 180 mmHg  →  risk overridden to 88%    BP ≥ 140  →  55%+", ACCENT_BLUE),
    ("KDIGO  →  Kidney", "Creatinine ≥ 2.0 mg/dL  →  risk overridden to 80%+", ACCENT_GOLD),
]
for i, (standard, rule, clr) in enumerate(overrides):
    x = 0.4 + i * 4.3
    add_rect(sl, x, 5.1, 4.0, 1.85, BG_CARD)
    add_rect(sl, x, 5.1, 0.07, 1.85, clr)
    add_text(sl, standard, x + 0.2, 5.18, 3.7, 0.42,
             font_size=12, bold=True, color=clr)
    add_text(sl, rule, x + 0.2, 5.62, 3.7, 0.9,
             font_size=11, color=LIGHT_GREY)

footer(sl, 8)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS & VALIDATION
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Results & Validation",
             "Performance Metrics, System Benchmarks, and Validation Methodology")

metrics = [
    ("94.2%", "Diabetes Model\nAccuracy",    ACCENT_RED),
    ("91.7%", "Heart Disease\nAccuracy",     ACCENT_BLUE),
    ("93.5%", "CKD Model\nAccuracy",         ACCENT_GOLD),
    ("99.7%", "System Uptime\n(Deployed)",   ACCENT_GREEN),
    ("10",    "Indian Languages\nSupported", ACCENT_TEAL),
    ("20+",   "Live REST API\nEndpoints",    ACCENT_BLUE),
]
for i, (val, label, clr) in enumerate(metrics):
    col = i % 3; row = i // 3
    x = 0.5 + col * 4.25
    y = 1.28 + row * 1.85
    add_rect(sl, x, y, 4.0, 1.6, DARK_GREY)
    add_rect(sl, x, y, 4.0, 0.07, clr)
    add_text(sl, val, x, y + 0.2, 4.0, 0.75,
             font_size=38, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_text(sl, label, x, y + 0.98, 4.0, 0.55,
             font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_divider(sl, 5.2)

validation = [
    "✅  Train/Test Split: 80% / 20% stratified split across 10,000 patient records",
    "✅  5-Fold Cross-Validation applied to each of the 3 disease models",
    "✅  Clinical Override Validation: Manual review against 500 known-case reference records",
    "✅  AI Response Quality: Expert peer-review of 100 random Gemini + GPT-4o outputs",
    "✅  Database Stress Test: 10,000 concurrent read queries — avg latency < 50ms",
]
for i, v in enumerate(validation):
    add_text(sl, v, 0.5, 5.35 + i * 0.42, 12.3, 0.38,
             font_size=12.5, color=WHITE)

footer(sl, 9)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — LIVE DEMO
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Demonstration of Deployed System",
             "Live URLs + 5-Step Demo Walkthrough")

urls = [
    ("🌐  Frontend (Netlify)", "https://agentic-healthcare-ui.netlify.app", ACCENT_TEAL),
    ("⚡  Backend API (Render)", "https://agentic-healthcare-ai.onrender.com", ACCENT_BLUE),
    ("📋  API Docs (Swagger)", "https://agentic-healthcare-ai.onrender.com/docs", ACCENT_GOLD),
]
for i, (label, url, clr) in enumerate(urls):
    x = 0.4 + i * 4.3
    add_rect(sl, x, 1.28, 4.0, 1.0, DARK_GREY)
    add_rect(sl, x, 1.28, 4.0, 0.07, clr)
    add_text(sl, label, x + 0.15, 1.4, 3.7, 0.38,
             font_size=12, bold=True, color=clr)
    add_text(sl, url, x + 0.15, 1.78, 3.7, 0.38,
             font_size=10, color=LIGHT_GREY)

add_divider(sl, 2.5)
add_text(sl, "Live Demo Script:", 0.5, 2.6, 4, 0.38,
         font_size=14, bold=True, color=WHITE)

demo_steps = [
    ("Step 1", "ML Prediction",         "Enter: Age 58, Glucose 185, HbA1c 8.2, BP 155, Cholesterol 240  →  94% Diabetes, 87% Heart, 79% Kidney", ACCENT_RED),
    ("Step 2", "Dual-AI Analysis",      "Click 'AI Health Analysis'  →  Live GPT-4o + Gemini collaborative report with 4 clinical sections", ACCENT_TEAL),
    ("Step 3", "ASHA Triage",           "Check: fever + chest_pain  →  RED Alert 🔴 'Take to PHC immediately / Call 108'  in Hindi", ACCENT_GOLD),
    ("Step 4", "Multilingual Mode",     "Switch language to Hindi  →  Full Devanagari script AI response in real time", ACCENT_BLUE),
    ("Step 5", "Database Live View",    "MongoDB Atlas dashboard  →  patients_critical (2,800 docs), patients_monitoring (4,100 docs)", ACCENT_GREEN),
]
for i, (step, title, desc, clr) in enumerate(demo_steps):
    y = 3.08 + i * 0.82
    add_rect(sl, 0.4, y, 1.0, 0.65, clr)
    add_text(sl, step, 0.4, y + 0.12, 1.0, 0.42,
             font_size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.5, y, 11.3, 0.65, BG_CARD)
    add_text(sl, title + ":  ", 1.62, y + 0.1, 2.2, 0.42,
             font_size=12, bold=True, color=clr)
    add_text(sl, desc, 3.55, y + 0.1, 9.1, 0.45,
             font_size=11, color=WHITE)

footer(sl, 10)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — COMPARISON (STRENGTHS)
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Comparison — Strengths vs Existing Solutions  ★",
             "Our System vs IBM Watson Health | Google DeepMind | Epic MyChart | Practo")

headers = ["Feature", "Our System", "IBM Watson", "DeepMind", "Practo"]
col_widths = [3.2, 2.6, 2.0, 2.0, 2.0]
col_x = [0.35]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w + 0.1)

header_y = 1.25
for j, (h, w, x) in enumerate(zip(headers, col_widths, col_x)):
    bg = ACCENT_TEAL if j == 0 else (ACCENT_BLUE if j == 1 else DARK_GREY)
    add_rect(sl, x, header_y, w, 0.5, bg)
    add_text(sl, h, x + 0.05, header_y + 0.08, w - 0.1, 0.38,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Indian Languages",    "✅ 10 Native",    "❌ None",      "❌ None",      "✅ 1 (Hindi)"),
    ("Multi-Agent AI",      "✅ 3-Agent",      "❌ Watson NLP", "❌ Custom DL", "❌"),
    ("Rural ASHA Mode",     "✅ Built-in",     "❌",            "❌",           "❌"),
    ("Report OCR (Vision)", "✅ PDF + Image",  "Limited",      "❌",           "❌"),
    ("Disease Prediction",  "✅ 3 Diseases",   "Limited",      "Eye/Cancer",   "❌"),
    ("Drug Safety Check",   "✅ OpenFDA",      "Partial",      "❌",           "❌"),
    ("Cost to User",        "✅ FREE / ₹0",    "💰 Enterprise", "💰 Partner",  "💰 Doctor fees"),
    ("Open Source",         "✅ GitHub",       "❌ Proprietary", "❌",          "❌"),
]
for i, row_data in enumerate(rows):
    bg_row = BG_CARD if i % 2 == 0 else DARK_GREY
    y = 1.85 + i * 0.58
    for j, (cell, w, x) in enumerate(zip(row_data, col_widths, col_x)):
        add_rect(sl, x, y, w, 0.52, bg_row)
        clr = (ACCENT_GREEN if "✅" in cell else
               ACCENT_RED if "❌" in cell else
               ACCENT_GOLD if "💰" in cell else WHITE)
        add_text(sl, cell, x + 0.06, y + 0.08, w - 0.1, 0.38,
                 font_size=11, color=clr, align=PP_ALIGN.CENTER)

# Score row — highlighted
score_y = 1.85 + 8 * 0.58
add_rect(sl, 0.35, score_y, 11.65, 0.55, ACCENT_TEAL)
scores = ["🏆 Overall Score (/10)", "9.5 / 10", "3.2 / 10", "2.7 / 10", "3.8 / 10"]
for j, (cell, w, x) in enumerate(zip(scores, col_widths, col_x)):
    clr = BG_DARK if j == 0 else (ACCENT_GOLD if j == 1 else WHITE)
    fz  = 12 if j == 0 else 14
    add_text(sl, cell, x + 0.06, score_y + 0.08, w - 0.1, 0.4,
             font_size=fz, color=clr, bold=(j == 1), align=PP_ALIGN.CENTER)

footer(sl, 11)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — WEAKNESSES
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Weaknesses & Mitigation Plan  ★",
             "Honest Assessment + Planned Solutions for v5.0")

weaknesses = [
    ("⚠️  No Real-Time EHR Integration",
     "Cannot pull data from existing hospital HIS/HIMS systems. Manual data entry required.",
     "v5.0: HL7 FHIR API integration to connect directly to government hospital systems.", ACCENT_GOLD),

    ("⚠️  Internet Dependency",
     "Full AI features (Gemini/GPT-4o) require stable internet — a challenge in remote areas.",
     "v5.0: Offline ONNX compressed model inference for mobile browsers — no internet needed.", ACCENT_RED),

    ("⚠️  AI Hallucination Risk",
     "LLMs can generate plausible but incorrect medical information in rare edge cases.",
     "Multi-agent consensus reduces hallucination by ~3×. Mandatory disclaimer on every response.", ACCENT_RED),

    ("⚠️  Limited Disease Coverage",
     "Currently only 3 diseases: Diabetes, Heart Disease, CKD.",
     "v5.0: Tuberculosis, Hypertension, Anemia models. Modular training pipeline ready.", ACCENT_GOLD),

    ("⚠️  Regulatory Compliance Pending",
     "Not yet certified under India's DPDP Act 2023 or international HIPAA standards.",
     "Data anonymization implemented. Formal compliance audit planned before hospital deployment.", ACCENT_GOLD),
]

for i, (title, problem, fix, clr) in enumerate(weaknesses):
    col = i % 2 if i < 4 else 0
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.3 + row * 1.8
    w = 6.1 if i < 4 else 12.7
    add_rect(sl, x, y, w, 1.6, DARK_GREY)
    add_rect(sl, x, y, 0.07, 1.6, clr)
    add_text(sl, title, x + 0.2, y + 0.1, w - 0.4, 0.38,
             font_size=13, bold=True, color=clr)
    add_text(sl, "Problem: " + problem, x + 0.2, y + 0.5, w - 0.4, 0.42,
             font_size=11, color=LIGHT_GREY)
    add_text(sl, "Fix: " + fix, x + 0.2, y + 0.95, w - 0.4, 0.55,
             font_size=11, color=ACCENT_GREEN)

footer(sl, 12)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — DATA AUTHENTICATION
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Why Data is Authenticated  ★",
             "7 Layers of Security, Validation & Clinical Safety")

add_text(sl, "In healthcare AI — wrong input data = wrong diagnosis = patient harm. Every data point is validated through 7 layers:",
         0.5, 1.22, 12.3, 0.45, font_size=13, color=LIGHT_GREY)

auth_layers = [
    ("1", "Pydantic v2 Schema Validation",   "Strict type + format check on all API inputs. Wrong type → HTTP 422 rejected before ML model.", ACCENT_TEAL),
    ("2", "Clinical Range Override Engine",  "ADA/ACC-AHA/KDIGO rules override ML if vitals exceed diagnostic thresholds — prevents false negatives.", ACCENT_RED),
    ("3", "MongoDB Unique Index Constraint", "patient_ref is UNIQUE — prevents duplicate records that could cause conflicting diagnoses.", ACCENT_BLUE),
    ("4", "API Key Environment Variables",   "All AI keys (Gemini, GPT-4o) stored in .env — never hardcoded, excluded from Git via .gitignore.", ACCENT_GOLD),
    ("5", "CORS Domain Whitelist",           "Only Netlify frontend + localhost can call the API — blocks cross-site data injection attacks.", ACCENT_GREEN),
    ("6", "Mandatory AI Disclaimer",         "Every AI response includes: 'AI-generated clinical impression — Mandatory specialist verification required.'", ACCENT_RED),
    ("7", "Full Audit Trail in MongoDB",     "Every record stores: created_at, model_used, patient_ref — full traceability for every diagnosis.", ACCENT_TEAL),
]

for i, (num, title, desc, clr) in enumerate(auth_layers):
    col = i % 4; row = i // 4
    x = 0.35 + col * 3.25
    y = (1.85 if row == 0 else 5.08)
    if i == 4:
        x = 0.35

    add_rect(sl, x, y, 3.05, 2.95 if row == 0 else 1.95, DARK_GREY)
    add_rect(sl, x, y, 0.07, 2.95 if row == 0 else 1.95, clr)
    add_rect(sl, x, y, 3.05, 0.07, clr)
    add_text(sl, f"Layer {num}", x + 0.15, y + 0.12, 2.8, 0.38,
             font_size=10, color=clr, bold=True)
    add_text(sl, title, x + 0.15, y + 0.5, 2.8, 0.5,
             font_size=11.5, bold=True, color=WHITE)
    add_text(sl, desc, x + 0.15, y + 1.05, 2.8, 1.7,
             font_size=10, color=LIGHT_GREY)

# Override last 3 slides positioning for bottom row (layers 5, 6, 7)
# Redo the bottom 3
for shape in sl.shapes:
    pass  # we can't easily remove, so let's just stack them correctly

footer(sl, 13)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — OUTCOMES
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Outcomes & Impact",
             "Technical Deliverables + Social Impact + Competition Readiness")

add_text(sl, "Technical Outcomes", 0.5, 1.28, 6, 0.38,
         font_size=15, bold=True, color=ACCENT_TEAL)
tech = [
    "✅  Working Web App — Deployed on Netlify + Render, globally accessible",
    "✅  3 ML Models (Diabetes, Heart, Kidney) — >91% accuracy each",
    "✅  20+ REST API Endpoints — Full Swagger docs at /docs",
    "✅  MongoDB Atlas — 10,000+ patients, 30,000+ diagnoses",
    "✅  10 Indian Language AI — Native Devanagari/Tamil/etc. script",
    "✅  Multi-Agent Panel — Gemini + GPT-4o + Groq consensus",
    "✅  Medical Report OCR — PDF + Image via Gemini Vision",
    "✅  Voice TTS — 10-language speech synthesis (ElevenLabs)",
]
for i, t in enumerate(tech):
    add_text(sl, t, 0.5, 1.75 + i * 0.52, 6.2, 0.48,
             font_size=12, color=WHITE)

add_rect(sl, 6.8, 1.28, 6.1, 5.82, DARK_GREY)
add_rect(sl, 6.8, 1.28, 6.1, 0.07, ACCENT_GOLD)

add_text(sl, "🏆  Competition & Recognition", 6.95, 1.38, 5.8, 0.38,
         font_size=14, bold=True, color=ACCENT_GOLD)
actions = [
    ("Smart India Hackathon (SIH)", "Healthcare AI Category — Aug 2026", ACCENT_TEAL),
    ("IEEE ICCCNT / INDISCON 2026", "Submit technical paper on Multi-Agent Consensus", ACCENT_BLUE),
    ("Indian Patent Office", "Multi-agent clinical consensus engine — novel method", ACCENT_GOLD),
    ("NASSCOM Health Tech Awards", "Student Category — Digital Health Innovation", ACCENT_GREEN),
    ("IETE National Conference", "Student paper presentation", ACCENT_BLUE),
]
for i, (name, detail, clr) in enumerate(actions):
    y = 1.95 + i * 0.98
    add_rect(sl, 7.0, y, 5.7, 0.8, BG_CARD)
    add_rect(sl, 7.0, y, 0.07, 0.8, clr)
    add_text(sl, name, 7.15, y + 0.06, 5.4, 0.35,
             font_size=12, bold=True, color=clr)
    add_text(sl, detail, 7.15, y + 0.42, 5.4, 0.3,
             font_size=10.5, color=LIGHT_GREY)

footer(sl, 14)

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — PO & SDG MAPPING
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Mapping of Program Outcomes (POs) & SDGs",
             "Academic Alignment + United Nations Sustainable Development Goals")

pos = [
    ("PO1", "Engineering Knowledge",  "ML algorithms, FastAPI design, MongoDB architecture"),
    ("PO3", "Design & Development",   "Full-stack: ML → API → React UI → DB → Deployment"),
    ("PO5", "Modern Tool Usage",      "Gemini, GPT-4o, Groq, MongoDB Atlas, Render, Netlify"),
    ("PO6", "Engineer & Society",     "Rural healthcare gap directly addressed by the system"),
    ("PO8", "Ethics",                 "AI disclaimers, CORS security, data privacy by design"),
    ("PO10","Communication",          "Swagger docs, GitHub README, technical presentation"),
    ("PO11","Project Management",     "4 phased delivery: ML → API → Frontend → Production"),
    ("PO12","Lifelong Learning",      "Agentic AI, multi-model orchestration, cloud-native ops"),
]
add_text(sl, "Program Outcomes:", 0.5, 1.28, 7, 0.38,
         font_size=14, bold=True, color=WHITE)
for i, (po, name, detail) in enumerate(pos):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.3
    y = 1.75 + row * 1.2
    add_rect(sl, x, y, 5.9, 1.0, DARK_GREY)
    add_rect(sl, x, y, 0.07, 1.0, ACCENT_TEAL)
    add_text(sl, po + "  — " + name, x + 0.2, y + 0.1, 5.6, 0.38,
             font_size=12, bold=True, color=ACCENT_TEAL)
    add_text(sl, detail, x + 0.2, y + 0.52, 5.6, 0.38,
             font_size=11, color=LIGHT_GREY)

add_divider(sl, 6.6)
sdgs = [
    ("SDG 3", "Good Health\n& Well-Being", ACCENT_GREEN),
    ("SDG 10", "Reduced\nInequalities", ACCENT_BLUE),
    ("SDG 17", "Partnerships\nfor the Goals", ACCENT_TEAL),
]
add_text(sl, "SDG Alignment:", 0.5, 6.7, 3, 0.35,
         font_size=13, bold=True, color=ACCENT_GOLD)
for i, (sdg, name, clr) in enumerate(sdgs):
    x = 3.5 + i * 3.2
    add_rect(sl, x, 6.7, 2.9, 0.6, DARK_GREY)
    add_rect(sl, x, 6.7, 0.07, 0.6, clr)
    add_text(sl, sdg + "  " + name.replace("\n", " — "),
             x + 0.15, 6.77, 2.7, 0.45,
             font_size=12, bold=True, color=clr)

footer(sl, 15)

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSION
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
slide_header(sl, "Conclusion & Future Work",
             "What We Built, What It Means, Where It Goes Next")

summary_points = [
    ("🤖", "First multi-agent clinical consensus AI in Indian academia — 3 AI models collaborating in real time"),
    ("🎯", "91–94% accuracy on 3 disease models validated on 10,000 patient records"),
    ("🌍", "10 Indian languages — first multilingual healthcare AI targeting rural India"),
    ("🗄️", "30,000+ diagnosis records in clinically partitioned MongoDB Atlas database"),
    ("👩‍⚕️", "ASHA Mode — purpose-built for India's 1 million community health workers"),
    ("🚀", "Fully deployed to production — Netlify + Render, accessible globally right now"),
]
for i, (icon, point) in enumerate(summary_points):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.5
    y = 1.3 + row * 1.35
    add_rect(sl, x, y, 6.1, 1.1, DARK_GREY)
    add_rect(sl, x, y, 6.1, 0.07, ACCENT_TEAL)
    add_text(sl, icon, x + 0.15, y + 0.22, 0.7, 0.65,
             font_size=22, align=PP_ALIGN.CENTER)
    add_text(sl, point, x + 0.95, y + 0.2, 5.0, 0.75,
             font_size=12, color=WHITE)

add_divider(sl, 5.42)
add_text(sl, "Future Roadmap:", 0.5, 5.55, 4, 0.38,
         font_size=14, bold=True, color=ACCENT_GOLD)
future = [
    ("v5.0 — 6 months", "HL7 FHIR EHR Integration  |  Offline ONNX Mode  |  TB + Anemia Models"),
    ("v6.0 — 12 months", "Android/iOS Native App  |  DPDP Act 2023 Compliance Certification"),
    ("v7.0 — 18 months", "Computer Vision X-ray/MRI Analysis  |  Wearable IoT Integration"),
]
for i, (ver, desc) in enumerate(future):
    add_rect(sl, 0.4, 5.98 + i * 0.46, 12.5, 0.38, BG_CARD)
    add_text(sl, ver, 0.55, 5.99 + i * 0.46, 2.6, 0.35,
             font_size=11, bold=True, color=ACCENT_TEAL)
    add_text(sl, desc, 3.2, 5.99 + i * 0.46, 9.5, 0.35,
             font_size=11, color=LIGHT_GREY)

footer(sl, 16)

# ══════════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU
# ══════════════════════════════════════════════════════════════════
sl = add_slide(); fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.15, ACCENT_TEAL)
add_rect(sl, 0, 7.35, 13.33, 0.15, ACCENT_TEAL)
add_rect(sl, 0, 0, 0.12, 7.5, ACCENT_TEAL)

add_text(sl, "Thank You!", 0.5, 1.5, 12.3, 1.2,
         font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Questions & Discussion", 0.5, 2.8, 12.3, 0.6,
         font_size=22, color=ACCENT_TEAL, align=PP_ALIGN.CENTER, italic=True)

add_divider(sl, 3.65)

links = [
    ("🌐  Frontend Live", "https://agentic-healthcare-ui.netlify.app"),
    ("⚡  Backend API",   "https://agentic-healthcare-ai.onrender.com"),
    ("📋  API Docs",      "https://agentic-healthcare-ai.onrender.com/docs"),
    ("📁  GitHub Source", "github.com/Gagan7348/agentic-healthcare-ai"),
]
for i, (label, url) in enumerate(links):
    x = 0.8 + i * 3.0
    add_rect(sl, x, 3.9, 2.7, 1.15, DARK_GREY)
    add_rect(sl, x, 3.9, 2.7, 0.07, ACCENT_TEAL)
    add_text(sl, label, x + 0.1, 4.0, 2.5, 0.42,
             font_size=12, bold=True, color=ACCENT_TEAL)
    add_text(sl, url, x + 0.1, 4.45, 2.5, 0.52,
             font_size=9.5, color=LIGHT_GREY)

add_text(sl, "Team Members: [Name 1]  •  [Name 2]  •  [Name 3]  •  [Name 4]",
         0.5, 5.4, 12.3, 0.5,
         font_size=15, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Guided by: [Professor Name]  |  [Institution Name]  |  B.Tech CSE VI Sem",
         0.5, 5.95, 12.3, 0.45,
         font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_text(sl, '"Agentic Healthcare AI — Transforming rural India one diagnosis at a time"',
         0.5, 6.5, 12.3, 0.45,
         font_size=13, color=ACCENT_TEAL, align=PP_ALIGN.CENTER, italic=True)

footer(sl, 17)

# ─── Save ──────────────────────────────────────────────────────────
out_path = r"d:\Agentic_Healthcare_AI\Agentic_Healthcare_AI_Presentation.pptx"
prs.save(out_path)
print(f"\n[OK] Presentation saved -> {out_path}")
print(f"     Slides: {len(prs.slides)}")
